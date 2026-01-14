#!/usr/bin/env python3
"""Analyze PPT structure from example files."""

import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from collections import defaultdict

PPT_DIR = "/Users/bisuv/Documents/internProject/pptgen_workflow/external_data/精品PPT-版式"

def analyze_ppt(filepath: str) -> dict:
    """Analyze a single PPT file and return structure info."""
    try:
        prs = Presentation(filepath)
    except Exception as e:
        return {"error": str(e), "filename": os.path.basename(filepath)}
    
    filename = os.path.basename(filepath)
    
    slides_info = []
    for idx, slide in enumerate(prs.slides):
        slide_data = {
            "index": idx + 1,
            "shapes_count": len(slide.shapes),
            "has_title": False,
            "title_text": "",
            "text_boxes": 0,
            "images": 0,
            "charts": 0,
            "tables": 0,
        }
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_data["text_boxes"] += 1
                if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                    if shape.placeholder_format.type == 1:  # Title
                        slide_data["has_title"] = True
                        slide_data["title_text"] = shape.text[:50] if shape.text else ""
            if shape.shape_type == 13:  # Picture
                slide_data["images"] += 1
            if shape.has_chart:
                slide_data["charts"] += 1
            if shape.has_table:
                slide_data["tables"] += 1
        
        # Try to get title from first shape if not found
        if not slide_data["has_title"] and slide.shapes:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    slide_data["title_text"] = shape.text.strip()[:50]
                    break
        
        slides_info.append(slide_data)
    
    # Categorize slides by position
    total_slides = len(slides_info)
    categories = {
        "cover": 0,  # First 1-2 slides
        "intro": 0,  # Next 1-2 slides (objectives, outline)
        "content": 0,  # Main content
        "summary": 0,  # Last 1-2 slides
    }
    
    if total_slides >= 1:
        categories["cover"] = 1
    if total_slides >= 3:
        categories["intro"] = min(2, total_slides - 2)
    if total_slides >= 4:
        categories["summary"] = min(2, total_slides - 3)
    categories["content"] = max(0, total_slides - categories["cover"] - categories["intro"] - categories["summary"])
    
    # Extract field from filename
    field = "unknown"
    if "【" in filename and "】" in filename:
        field = filename.split("【")[1].split("】")[0]
    
    return {
        "filename": filename,
        "field": field,
        "total_slides": total_slides,
        "structure": categories,
        "slides_detail": slides_info[:5],  # First 5 slides detail
        "avg_images_per_slide": sum(s["images"] for s in slides_info) / max(1, total_slides),
        "avg_textboxes_per_slide": sum(s["text_boxes"] for s in slides_info) / max(1, total_slides),
    }


def main():
    results = []
    
    for filename in sorted(os.listdir(PPT_DIR)):
        if filename.endswith(".pptx") and not filename.startswith(".~"):
            filepath = os.path.join(PPT_DIR, filename)
            print(f"Analyzing: {filename}")
            result = analyze_ppt(filepath)
            results.append(result)
            print(f"  -> {result.get('total_slides', 'error')} slides, field: {result.get('field', 'unknown')}")
    
    # Summary statistics
    print("\n" + "=" * 60)
    print("SUMMARY STATISTICS")
    print("=" * 60)
    
    valid_results = [r for r in results if "error" not in r]
    
    if valid_results:
        slide_counts = [r["total_slides"] for r in valid_results]
        print(f"Total PPTs analyzed: {len(valid_results)}")
        print(f"Slide count range: {min(slide_counts)} - {max(slide_counts)}")
        print(f"Average slides: {sum(slide_counts) / len(slide_counts):.1f}")
        print(f"Median slides: {sorted(slide_counts)[len(slide_counts)//2]}")
        
        # Group by field
        by_field = defaultdict(list)
        for r in valid_results:
            by_field[r["field"]].append(r["total_slides"])
        
        print("\nBy Field:")
        for field, counts in sorted(by_field.items()):
            print(f"  {field}: avg={sum(counts)/len(counts):.1f}, range={min(counts)}-{max(counts)}, count={len(counts)}")
    
    # Output JSON for further analysis
    output_path = os.path.join(os.path.dirname(PPT_DIR), "ppt_analysis_results.json")
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(results, f, ensure_ascii=False, indent=2)
    print(f"\nDetailed results saved to: {output_path}")


if __name__ == "__main__":
    main()
